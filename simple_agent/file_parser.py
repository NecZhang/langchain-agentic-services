"""Utilities for loading and converting different file types to plain text.

This module provides helper functions to extract text from a variety of
common file formats.  In order to keep the code lightweight and self
contained, it relies only on packages that are preinstalled in the
execution environment.  Currently supported formats include:

* Plain text files (``.txt``)
* JSON files (``.json``) – the contents are loaded and pretty‑printed
  back to a JSON string
* Portable Document Format files (``.pdf``) via PyMuPDF (``fitz``)
* PowerPoint presentations (``.pptx``) via ``python‑pptx``

Other formats like Microsoft Word documents (``.doc``/``.docx``) and
images (``.jpg``/``.png``) would normally require external libraries
(``python‑docx``, ``pytesseract``, etc.), which are not installed in
this environment.  Attempting to parse unsupported formats will
trigger a ``NotImplementedError``.
"""

from __future__ import annotations

import json
import os

import fitz  # PyMuPDF is available in this environment
from pptx import Presentation

__all__ = ["parse_file"]


def _parse_text_file(path: str) -> str:
    """Read a plain text file into a string.

    Parameters
    ----------
    path : str
        Path to the file on disk.

    Returns
    -------
    str
        The contents of the file decoded as UTF‑8.  Any decoding
        errors are ignored to maximise robustness.
    """
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def _parse_json_file(path: str) -> str:
    """Load a JSON file and serialize it back to formatted JSON.

    Parameters
    ----------
    path : str
        Path to the file on disk.

    Returns
    -------
    str
        A pretty printed JSON string representing the loaded object.
    """
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        data = json.load(f)
    # Dump with indentation and ensure unicode is preserved
    return json.dumps(data, indent=2, ensure_ascii=False)


def _parse_pdf_file(path: str) -> str:
    """Extract text from a PDF using PyMuPDF.

    This function iterates through each page of the PDF and collects
    the text extracted by PyMuPDF.  While this approach works well
    for digitally generated PDFs, scanned documents will require
    optical character recognition (OCR), which is beyond the scope of
    this environment.  The returned string concatenates the text from
    all pages separated by a blank line.

    Parameters
    ----------
    path : str
        Path to the PDF file.

    Returns
    -------
    str
        The extracted text.
    """
    doc = fitz.open(path)
    texts = []
    for page_number, page in enumerate(doc):
        try:
            text = page.get_text("text")
        except Exception:
            text = ""
        if text:
            texts.append(text.strip())
    return "\n\n".join(texts)


def _parse_pptx_file(path: str) -> str:
    """Extract text from a PowerPoint (``.pptx``) file.

    The ``python‑pptx`` library is used to iterate through slides and
    shapes.  Only textual content from shapes is collected.  Images
    embedded in the presentation will not be processed.

    Parameters
    ----------
    path : str
        Path to the presentation file.

    Returns
    -------
    str
        A newline‑separated string of all text found in the slides.
    """
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            # many shape types expose a .text attribute (TextFrame)
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    parts.append(text)
    return "\n\n".join(parts)


def _parse_docx_file(path: str) -> str:
    """Extract text from a Word document (``.docx``) file.

    The ``python-docx`` library is used to iterate through paragraphs
    and extract text content. Tables and other elements are also processed.

    Parameters
    ----------
    path : str
        Path to the Word document file.

    Returns
    -------
    str
        A newline-separated string of all text found in the document.
    """
    try:
        from docx import Document
    except ImportError:
        raise NotImplementedError("DOC/DOCX parsing requires python-docx to be installed.")

    doc = Document(path)
    parts = []

    # Extract text from paragraphs
    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)

    # Extract text from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                if cell_text:
                    row_text.append(cell_text)
            if row_text:
                parts.append(" | ".join(row_text))

    return "\n\n".join(parts)


def _parse_image_file(path: str) -> str:
    """Extract text from an image file using OCR.

    The ``pytesseract`` library is used to perform optical character
    recognition on the image. The Tesseract engine must be installed
    on the system.

    Parameters
    ----------
    path : str
        Path to the image file.

    Returns
    -------
    str
        The extracted text from the image.
    """
    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        raise NotImplementedError("Image OCR requires pytesseract and Pillow to be installed.")

    try:
        # Open the image and extract text
        image = Image.open(path)
        text = pytesseract.image_to_string(image)
        return text.strip()
    except Exception as e:
        raise RuntimeError(f"Failed to extract text from image: {str(e)}")


def parse_file(path: str) -> str:
    """Parse a file into plain text.

    The parser dispatches based on the file extension.  Unsupported
    formats will raise a ``NotImplementedError``.  See module
    docstring for the list of supported formats.

    Parameters
    ----------
    path : str
        Location of the file on disk.

    Returns
    -------
    str
        Extracted textual contents from the file.

    Raises
    ------
    NotImplementedError
        If parsing for the given extension is not available.
    ValueError
        If the file extension is unknown.
    """
    if not os.path.exists(path):
        raise FileNotFoundError(f"File not found: {path}")
    ext = os.path.splitext(path)[1].lower()
    if ext in {".txt"}:
        return _parse_text_file(path)
    elif ext in {".json"}:
        return _parse_json_file(path)
    elif ext in {".pdf"}:
        return _parse_pdf_file(path)
    elif ext in {".pptx"}:
        return _parse_pptx_file(path)
    elif ext == ".docx":
        return _parse_docx_file(path)
    elif ext == ".doc":
        # Legacy .doc format is not supported by python-docx
        # Would require additional tools like python-docx2txt or LibreOffice
        raise NotImplementedError(
            "Legacy .doc format is not supported. Please convert to .docx format or use LibreOffice to convert the file."
        )
    elif ext in {".jpg", ".jpeg", ".png", ".bmp"}:
        return _parse_image_file(path)
    else:
        raise ValueError(f"Unsupported file extension: {ext}")
